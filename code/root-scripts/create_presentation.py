#!/usr/bin/env python3
"""
Create a 15-slide PowerPoint presentation about galaxy rotation curves and the causal-response model.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(2))
    title_frame = title_box.text_frame
    title_frame.text = "Galaxy Rotation Curves:\nThe Dark Matter Problem and a New Causal-Response Model"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.alignment = PP_ALIGN.CENTER
    title_para.font.color.rgb = RGBColor(0, 51, 102)
    
    author_box = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(8), Inches(1))
    author_frame = author_box.text_frame
    author_frame.text = "Jonathan Washburn & Elshad Allahyarov\nRecognition Science Institute"
    author_para = author_frame.paragraphs[0]
    author_para.font.size = Pt(18)
    author_para.alignment = PP_ALIGN.CENTER
    
    # Slide 2: The Problem - Galaxy Rotation Curves
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "The Galaxy Rotation Curve Problem"
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "What we observe:"
    p = tf.add_paragraph()
    p.text = "• Stars in galaxies orbit faster than expected"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "• Velocity remains constant with radius (flat rotation curves)"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "• Newtonian gravity predicts v ∝ r^(-1/2) (Keplerian decline)"
    p.level = 0
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "The discrepancy:"
    p.font.bold = True
    p = tf.add_paragraph()
    p.text = "• Factor of 5-10× more gravitational force needed"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "• Observed in ALL galaxies for 50+ years"
    p.level = 0
    
    # Figure placeholder (PDF conversion needed)
    # pic_path = "frg_rotation_curves.pdf"
    
    # Slide 3: Two Competing Explanations
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Two Radical Solutions"
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "1. DARK MATTER (Standard Model)"
    p = tf.add_paragraph()
    p.text = "• 85% of matter is invisible particles"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Only interacts gravitationally"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Despite decades of searches: NO detection"
    p.level = 1
    p.font.color.rgb = RGBColor(255, 0, 0)
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "2. MODIFIED GRAVITY"
    p = tf.add_paragraph()
    p.text = "• Einstein's General Relativity breaks down"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• New physics at galactic scales"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• MOND: Acceleration scale a₀ ≈ 1.2×10⁻¹⁰ m/s²"
    p.level = 1
    
    # Slide 4: Dark Matter Paradigm
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Dark Matter: The Standard Paradigm"
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Successes:"
    p = tf.add_paragraph()
    p.text = "✓ Cosmic Microwave Background"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "✓ Large-scale structure formation"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "✓ Gravitational lensing"
    p.level = 0
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "Problems:"
    p.font.color.rgb = RGBColor(255, 0, 0)
    p = tf.add_paragraph()
    p.text = "✗ Cusp-core problem"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "✗ Missing satellites problem"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "✗ Too-big-to-fail problem"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "✗ Radial Acceleration Relation (tight baryon-dynamics link)"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "✗ NO particle detected (XENON, LUX, LHC all null)"
    p.level = 0
    p.font.bold = True
    
    # Slide 5: MOND Theory
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "MOND: Modified Newtonian Dynamics"
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Key Idea (Milgrom 1983):"
    p = tf.add_paragraph()
    p.text = "Gravity is modified below acceleration scale a₀"
    p.level = 0
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "Successes:"
    p = tf.add_paragraph()
    p.text = "✓ Fits rotation curves with NO free parameters"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "✓ Predicts Baryonic Tully-Fisher Relation"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "✓ Explains Radial Acceleration Relation"
    p.level = 0
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "Problems:"
    p.font.color.rgb = RGBColor(255, 0, 0)
    p = tf.add_paragraph()
    p.text = "✗ No compelling theoretical foundation"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "✗ Struggles with galaxy clusters"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "✗ Requires ad-hoc extensions for cosmology"
    p.level = 0
    
    # Slide 6: The Third Way - Causal Response
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "A Third Way: Causal-Response Model"
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Our Proposal:"
    p = tf.add_paragraph()
    p.text = "Gravitational memory effect at long timescales"
    p.level = 0
    p.font.bold = True
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "Key Features:"
    p = tf.add_paragraph()
    p.text = "• Phenomenological (like Fermi's weak interaction theory)"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "• Causal (respects special relativity)"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "• Conservative (admits thermodynamic realization)"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "• Falsifiable (concrete testable predictions)"
    p.level = 0
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "Mathematical Form:"
    p = tf.add_paragraph()
    p.text = "a_eff(r) = w(r) × a_baryon(r)"
    p.level = 0
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(0, 102, 204)
    
    # Slide 7: Linear Response Theory
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Mathematical Framework: Linear Response"
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Memory Kernel:"
    p = tf.add_paragraph()
    p.text = "Γ(τ) = (w-1)/τ★ × exp(-τ/τ★)"
    p.level = 0
    p.font.size = Pt(18)
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "Transfer Function:"
    p = tf.add_paragraph()
    p.text = "H(iω) = 1 + (w-1)/(1 + iωτ★)"
    p.level = 0
    p.font.size = Pt(18)
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "Physical Constraints:"
    p = tf.add_paragraph()
    p.text = "✓ Causality (no acausal response)"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "✓ Reality (real time-domain response)"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "✓ Fast-limit recovery (H→1 as ω→∞)"
    p.level = 0
    
    # Figure placeholder
    # pic_path = "frg_memory_kernel.pdf"
    
    # Slide 8: Model Parameters
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Global Parameters (7 total)"
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Fitted to 99 high-quality (Q=1) SPARC galaxies:"
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "1. α = 0.18 (time scaling exponent)"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "2. a₀ = 1.95×10⁻¹⁰ m/s² (acceleration scale)"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "3. r₀ = 17.79 kpc (radial scale)"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "4-7. Morphology parameters (gas fraction, surface brightness)"
    p.level = 0
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "Key Constraint:"
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 0, 0)
    p = tf.add_paragraph()
    p.text = "NO per-galaxy tuning allowed!"
    p.level = 0
    p.font.bold = True
    p.font.size = Pt(20)
    
    # Slide 9: Results - Performance
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Results: Outstanding Performance"
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Causal-Response Model:"
    p = tf.add_paragraph()
    p.text = "• Median χ²/N = 1.19"
    p.level = 0
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(0, 153, 0)
    p.font.bold = True
    p = tf.add_paragraph()
    p.text = "• RMS velocity residual = 17.3 km/s"
    p.level = 0
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "MOND (same constraints):"
    p = tf.add_paragraph()
    p.text = "• Median χ²/N = 1.79"
    p.level = 0
    p.font.size = Pt(20)
    p = tf.add_paragraph()
    p.text = "• RMS velocity residual = 21.5 km/s"
    p.level = 0
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "Improvement:"
    p.font.bold = True
    p = tf.add_paragraph()
    p.text = "33% reduction in median residual"
    p.level = 0
    p.font.size = Pt(22)
    p.font.color.rgb = RGBColor(255, 0, 0)
    p.font.bold = True
    p = tf.add_paragraph()
    p.text = "19.5% reduction in RMS velocity error"
    p.level = 0
    p.font.size = Pt(22)
    p.font.color.rgb = RGBColor(255, 0, 0)
    p.font.bold = True
    
    # Slide 10: Chi-squared Distribution
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = "χ²/N Distribution: Model vs MOND"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    
    # Figure: frg_chi2_dist_q1q2.pdf (add manually)
    
    # Slide 11: Dwarf vs Spiral Prediction
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = "Natural Prediction: Dwarfs vs Spirals"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    
    text_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(4), Inches(2))
    tf = text_box.text_frame
    tf.text = "Model predicts:\n~1.8× stronger modification in dwarfs\n\nObserved:\n1.77 ± 0.29\n\nPerfect agreement!"
    for para in tf.paragraphs:
        para.font.size = Pt(18)
    
    # Figure: frg_dwarf_spiral_q1q2.pdf (add manually)
    
    # Slide 12: Radial Acceleration Relation
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = "Radial Acceleration Relation (RAR)"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    
    # Figure: frg_rar_q1q2.pdf (add manually)
    
    # Slide 13: Falsification Tests
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Falsification Tests"
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Concrete, testable predictions:"
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "1. Galaxy Cluster Lensing:"
    p.font.bold = True
    p = tf.add_paragraph()
    p.text = "κ/κ_GR = 1.8 ± 0.3 at R = 20-50 kpc"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Testable with Euclid Space Telescope (2025-2030)"
    p.level = 1
    p.font.color.rgb = RGBColor(255, 0, 0)
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "2. Laboratory Gravity:"
    p.font.bold = True
    p = tf.add_paragraph()
    p.text = "Inverse-square law holds to |β| < 10⁻⁴"
    p.level = 1
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "3. Pulsar Timing:"
    p.font.bold = True
    p = tf.add_paragraph()
    p.text = "No deviations expected (fast timescales)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "4. Cosmology:"
    p.font.bold = True
    p = tf.add_paragraph()
    p.text = "ISW, CMB lensing predictions (requires relativistic extension)"
    p.level = 1
    
    # Slide 14: Comparison with Other Theories
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Comparison: Major Theories"
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Theory Comparison:"
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "ΛCDM (NFW): Good clusters/cosmology, requires tuning"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "MOND: Excellent rotation curves, poor clusters"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "TeVeS: Excellent rotation curves, marginal cosmology"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "f(R) gravity: Poor rotation curves, good cosmology"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "Verlinde (2017): Poor on all scales"
    p.level = 0
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "Causal-Response (this work):"
    p.font.bold = True
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(0, 102, 204)
    p = tf.add_paragraph()
    p.text = "✓ Good rotation curves (χ²/N = 1.19)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "✓ Zero per-galaxy parameters"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "✓ Falsifiable predictions"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "? Clusters/cosmology untested (needs relativistic extension)"
    p.level = 1
    
    # Slide 15: Conclusions
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Conclusions & Future Directions"
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Key Findings:"
    p = tf.add_paragraph()
    p.text = "• Causal-response model fits 99 galaxies with 7 global parameters"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "• 33% better than MOND (median χ²/N: 1.19 vs 1.79)"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "• Natural dwarf/spiral prediction (1.8× observed: 1.77±0.29)"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "• Eliminates gas-fraction residual bias"
    p.level = 0
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "Paradigm Shift:"
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 0, 0)
    p = tf.add_paragraph()
    p.text = "\"Dark matter\" at galactic scales may be gravitational memory, not particles"
    p.level = 0
    p.font.size = Pt(18)
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "Next Steps:"
    p.font.bold = True
    p = tf.add_paragraph()
    p.text = "• Euclid cluster lensing (2025-2030): Definitive test"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "• Relativistic extension for cosmology"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "• Theoretical foundation (quantum gravity? emergent?)"
    p.level = 0
    
    # Save presentation
    prs.save('aaaa.pptx')
    print("OK Presentation created: aaaa.pptx (15 slides)")
    print("\nSlide breakdown:")
    print("  1. Title slide")
    print("  2. The galaxy rotation curve problem")
    print("  3. Two competing explanations")
    print("  4. Dark matter paradigm")
    print("  5. MOND theory")
    print("  6. The third way: Causal-response model")
    print("  7. Linear response theory")
    print("  8. Model parameters")
    print("  9. Results: Outstanding performance")
    print(" 10. Chi-squared distribution")
    print(" 11. Dwarf vs spiral prediction")
    print(" 12. Radial acceleration relation")
    print(" 13. Falsification tests")
    print(" 14. Comparison with other theories")
    print(" 15. Conclusions & future directions")

if __name__ == "__main__":
    create_presentation()

