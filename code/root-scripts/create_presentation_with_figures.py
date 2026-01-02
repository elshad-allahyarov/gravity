#!/usr/bin/env python3
"""
Create presentation with actual figures embedded.
Convert PDF figures to PNG and embed them.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pdf2image import convert_from_path
import os

def pdf_to_png(pdf_path, output_path):
    """Convert PDF to PNG"""
    try:
        images = convert_from_path(pdf_path, dpi=300, first_page=1, last_page=1)
        if images:
            images[0].save(output_path, 'PNG')
            return True
    except Exception as e:
        print(f"Error converting {pdf_path}: {e}")
        return False

def create_presentation():
    # Convert PDFs to PNGs first
    print("Converting PDF figures to PNG...")
    figures_to_convert = [
        'frg_rotation_curves.pdf',
        'frg_memory_kernel.pdf',
        'frg_chi2_dist_q1q2.pdf',
        'frg_dwarf_spiral_q1q2.pdf',
        'frg_rar_q1q2.pdf'
    ]
    
    png_files = {}
    for pdf_file in figures_to_convert:
        if os.path.exists(pdf_file):
            png_file = pdf_file.replace('.pdf', '.png')
            if pdf_to_png(pdf_file, png_file):
                png_files[pdf_file] = png_file
                print(f"  OK {pdf_file} -> {png_file}")
            else:
                print(f"  FAILED {pdf_file}")
        else:
            print(f"  MISSING {pdf_file}")
    
    print(f"\nConverted {len(png_files)} figures")
    print("\nCreating presentation...")
    
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Title
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(2))
    tf = title_box.text_frame
    tf.text = "Galaxy Rotation Curves:\nThe Dark Matter Problem and a New Causal-Response Model"
    p = tf.paragraphs[0]
    p.font.size = Pt(36)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    p.font.color.rgb = RGBColor(0, 51, 102)
    
    author_box = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(8), Inches(1))
    tf2 = author_box.text_frame
    tf2.text = "Jonathan Washburn & Elshad Allahyarov\nRecognition Science Institute"
    p2 = tf2.paragraphs[0]
    p2.font.size = Pt(18)
    p2.alignment = PP_ALIGN.CENTER
    
    # Slide 2: Problem
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "The Galaxy Rotation Curve Problem"
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "What we observe:"
    p = tf.add_paragraph()
    p.text = "Stars orbit faster than expected from visible matter"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "Velocity stays constant with radius (flat rotation curves)"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "Newtonian gravity predicts v proportional to r^(-1/2)"
    p.level = 0
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "The discrepancy: 5-10x more gravitational force needed"
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 0, 0)
    
    # Add figure if available
    if 'frg_rotation_curves.pdf' in png_files:
        left = Inches(5.5)
        top = Inches(2)
        width = Inches(4)
        slide.shapes.add_picture(png_files['frg_rotation_curves.pdf'], left, top, width=width)
    
    # Slide 3: Two Solutions
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Two Radical Solutions"
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "1. DARK MATTER"
    p = tf.add_paragraph()
    p.text = "85% of matter is invisible particles"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Despite decades of searches: NO detection"
    p.level = 1
    p.font.color.rgb = RGBColor(255, 0, 0)
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "2. MODIFIED GRAVITY"
    p = tf.add_paragraph()
    p.text = "Einstein's GR breaks down at galactic scales"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "MOND: acceleration scale a0 = 1.2x10^-10 m/s^2"
    p.level = 1
    
    # Slide 4: Dark Matter
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Dark Matter: Successes & Problems"
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Successes:"
    p = tf.add_paragraph()
    p.text = "Cosmic Microwave Background"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "Large-scale structure"
    p.level = 0
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "Problems:"
    p.font.color.rgb = RGBColor(255, 0, 0)
    p = tf.add_paragraph()
    p.text = "Cusp-core problem"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "Missing satellites"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "Radial Acceleration Relation (tight baryon-dynamics link)"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "NO particle detected after 40 years"
    p.level = 0
    p.font.bold = True
    
    # Slide 5: MOND
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "MOND: Modified Newtonian Dynamics"
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Successes:"
    p = tf.add_paragraph()
    p.text = "Fits rotation curves with NO free parameters"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "Predicts Baryonic Tully-Fisher Relation"
    p.level = 0
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "Problems:"
    p.font.color.rgb = RGBColor(255, 0, 0)
    p = tf.add_paragraph()
    p.text = "No theoretical foundation"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "Struggles with galaxy clusters"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "Requires ad-hoc extensions for cosmology"
    p.level = 0
    
    # Slide 6: Causal Response
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "A Third Way: Causal-Response Model"
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Gravitational memory effect at long timescales"
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "Key Features:"
    p = tf.add_paragraph()
    p.text = "Phenomenological (like Fermi's weak interaction)"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "Causal (respects special relativity)"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "Falsifiable (concrete testable predictions)"
    p.level = 0
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "Mathematical Form:"
    p = tf.add_paragraph()
    p.text = "a_eff(r) = w(r) x a_baryon(r)"
    p.level = 0
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(0, 102, 204)
    p.font.bold = True
    
    # Slide 7: Linear Response
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    tf = title_box.text_frame
    tf.text = "Mathematical Framework: Linear Response"
    p = tf.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    
    if 'frg_memory_kernel.pdf' in png_files:
        left = Inches(0.5)
        top = Inches(1.2)
        width = Inches(9)
        slide.shapes.add_picture(png_files['frg_memory_kernel.pdf'], left, top, width=width)
    
    # Slide 8: Parameters
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "7 Global Parameters (NO per-galaxy tuning)"
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Fitted to 99 high-quality SPARC galaxies:"
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "alpha = 0.18 (time scaling)"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "a0 = 1.95x10^-10 m/s^2 (acceleration scale)"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "r0 = 17.79 kpc (radial scale)"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "+ 4 morphology parameters"
    p.level = 0
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "Key: NO per-galaxy tuning!"
    p.font.bold = True
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(255, 0, 0)
    
    # Slide 9: Results
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Results: Outstanding Performance"
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Causal-Response Model:"
    p = tf.add_paragraph()
    p.text = "Median chi^2/N = 1.19"
    p.level = 0
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(0, 153, 0)
    p.font.bold = True
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "MOND (same constraints):"
    p = tf.add_paragraph()
    p.text = "Median chi^2/N = 1.79"
    p.level = 0
    p.font.size = Pt(24)
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "33% reduction in median residual"
    p.font.size = Pt(26)
    p.font.color.rgb = RGBColor(255, 0, 0)
    p.font.bold = True
    
    # Slide 10: Chi2 Distribution
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    tf = title_box.text_frame
    tf.text = "Chi-squared Distribution: Model vs MOND"
    p = tf.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    
    if 'frg_chi2_dist_q1q2.pdf' in png_files:
        left = Inches(0.5)
        top = Inches(1.2)
        width = Inches(9)
        slide.shapes.add_picture(png_files['frg_chi2_dist_q1q2.pdf'], left, top, width=width)
    
    # Slide 11: Dwarf vs Spiral
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    tf = title_box.text_frame
    tf.text = "Natural Prediction: Dwarfs vs Spirals"
    p = tf.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    
    text_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(4), Inches(2))
    tf = text_box.text_frame
    tf.text = "Predicted: 1.8x stronger in dwarfs\n\nObserved: 1.77 +/- 0.29\n\nPerfect agreement!"
    for para in tf.paragraphs:
        para.font.size = Pt(20)
        para.font.bold = True
    
    if 'frg_dwarf_spiral_q1q2.pdf' in png_files:
        left = Inches(5)
        top = Inches(1.2)
        width = Inches(4.5)
        slide.shapes.add_picture(png_files['frg_dwarf_spiral_q1q2.pdf'], left, top, width=width)
    
    # Slide 12: RAR
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    tf = title_box.text_frame
    tf.text = "Radial Acceleration Relation (RAR)"
    p = tf.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    
    if 'frg_rar_q1q2.pdf' in png_files:
        left = Inches(0.5)
        top = Inches(1.2)
        width = Inches(9)
        slide.shapes.add_picture(png_files['frg_rar_q1q2.pdf'], left, top, width=width)
    
    # Slide 13: Falsification
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Falsification Tests"
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "1. Galaxy Cluster Lensing:"
    p = tf.add_paragraph()
    p.text = "kappa/kappa_GR = 1.8 +/- 0.3"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Testable with Euclid (2025-2030)"
    p.level = 1
    p.font.color.rgb = RGBColor(255, 0, 0)
    p.font.bold = True
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "2. Laboratory Gravity:"
    p = tf.add_paragraph()
    p.text = "Inverse-square law holds"
    p.level = 1
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "3. Pulsar Timing:"
    p = tf.add_paragraph()
    p.text = "No deviations (fast timescales)"
    p.level = 1
    
    # Slide 14: Comparison
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Comparison with Other Theories"
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "LCDM: Good clusters, requires tuning"
    p = tf.add_paragraph()
    p.text = "MOND: Excellent rotation curves, poor clusters"
    p = tf.add_paragraph()
    p.text = "TeVeS: Good rotation curves, marginal cosmology"
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "Causal-Response (this work):"
    p.font.bold = True
    p.font.size = Pt(22)
    p.font.color.rgb = RGBColor(0, 102, 204)
    p = tf.add_paragraph()
    p.text = "Good rotation curves (chi^2/N = 1.19)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Zero per-galaxy parameters"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Falsifiable predictions"
    p.level = 1
    
    # Slide 15: Conclusions
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Conclusions"
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Key Findings:"
    p = tf.add_paragraph()
    p.text = "7 global parameters fit 99 galaxies"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "33% better than MOND"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "Natural dwarf/spiral prediction"
    p.level = 0
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "Paradigm Shift:"
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 0, 0)
    p = tf.add_paragraph()
    p.text = "Dark matter at galactic scales may be gravitational memory, not particles"
    p.level = 0
    p.font.size = Pt(18)
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "Next: Euclid cluster lensing test (2025-2030)"
    p.font.bold = True
    
    # Save
    prs.save('aaaa.pptx')
    print("\nOK Presentation created: aaaa.pptx (15 slides with figures)")

if __name__ == "__main__":
    create_presentation()


